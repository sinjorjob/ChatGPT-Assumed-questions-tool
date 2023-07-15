# -*- coding:utf-8 -*-
from __future__ import annotations
import logging
from modules.presets import *
import os
import datetime
import json
import time
from modules.presets import *
from pypinyin import lazy_pinyin
from pptx import Presentation


def save_proxy(proxy):
    data = {}
    if os.path.exists(KEY_FILE):
        with open(KEY_FILE, "r") as f:
            data = json.load(f)
    data["proxy"] = proxy
    with open(KEY_FILE, "w") as f:
        json.dump(data, f, indent=4)
        
        
def save_key(api_key):
    data = {}
    if os.path.exists(KEY_FILE):
        with open(KEY_FILE, "r") as f:
            data = json.load(f)
    data["api_key"] = api_key
    with open(KEY_FILE, "w") as f:
        json.dump(data, f, indent=4)

def load_key():

    if not os.path.isfile(KEY_FILE):
        return None
    with open(KEY_FILE, "r") as f:
        data = json.load(f)
    return data.get("api_key")

def load_proxy():
    if not os.path.isfile(KEY_FILE):
        return None
    with open(KEY_FILE, "r") as f:
        data = json.load(f)
    return data.get("proxy")

def reset_textbox():
    logging.debug("テキストボックスをリセットする")
    return gr.update(value="")

def start_outputing():
    # キャンセルボタンを表示し、実行ボタンを非表示にする。
    return (
        gr.Button.update(visible=False), 
        gr.Button.update(visible=True)
    )

def end_outputing():
    # 実行ボタンを表示し、キャンセルボタンを非表示にする。
    return (
        gr.Button.update(visible=True),
        gr.Button.update(visible=False),
    )
def transfer_input(inputs):
    # 一回の応答で処理を終え、遅延を低減する
    #textbox = reset_textbox()
    #outputing = start_outputing()
    return (
        inputs,
        gr.update(value=""),
        gr.Button.update(visible=False),
        gr.Button.update(visible=True),
    )
    
      
        
def start_outputing():
    logging.debug("Cancelボタンを表示し、Sendボタンを非表示にする。")
    return gr.Button.update(visible=False), gr.Button.update(visible=True)
    
def construct_text(role, text):
    return {"role": role, "content": text}

def construct_user(text):
    return construct_text("user", text)

def construct_system(text):
    return construct_text("system", text)

def construct_assistant(text):
    return construct_text("assistant", text)
 
def reset(current_model, *args):
    """質問をリセットする"""
    return current_model.reset(*args)

def extract_pptx_content(pptx):
   
    # PowerPointの内容を取得する
    presentation = Presentation(pptx.name)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
    
    contexts = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    text = paragraph.text
                    contexts.append(text)
            elif shape.has_table:
                table_data = []
                table = shape.table
                table_data.append(table)
                for table in table_data:
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text
                            contexts.append(cell_text)
                
    
    output = '\n'.join(contexts)
    return output


def predict(current_model, uploadFile, favorable_question, skeptical_question, important_question, chatbot):
    """
    OpenAIのChatCompletion.create()メソッドを使用して、Chatbotの回答を取得します。
    """
    iter = current_model.predict(uploadFile, favorable_question, skeptical_question, important_question, chatbot)
    for chatbot in iter:
        time.sleep(0.03)
        yield chatbot

def set_temperature(current_model, *args):
    current_model.set_temperature(*args)
    
def set_top_p(current_model, *args):
    current_model.set_top_p(*args)
    
def set_stop_sequence(current_model, *args):
    current_model.set_stop_sequence(*args)
    
def set_presence_penalty(current_model, *args):
    current_model.set_presence_penalty(*args)

def set_frequency_penalty(current_model, *args):
    current_model.set_frequency_penalty(*args)
       
    
def set_key(current_model, *args):
    return current_model.set_key(*args)


    
def set_model(current_model, *args):
    return current_model.set_model(*args)

def hide_middle_chars(s):
    if s is None:
        return ""
    if len(s) <= 8:
        return s
    else:
        head = s[:4]
        tail = s[-4:]
        hidden = "*" * (len(s) - 8)
        return head + hidden + tail
    
    
def get_last_day_of_month(any_day):
    # 28日という日は、どの月にも存在する。4日後、それは常に来月である
    next_month = any_day.replace(day=28) + datetime.timedelta(days=4)
    # 現日数を引くと1ヶ月前に戻る
    return next_month - datetime.timedelta(days=next_month.day)

def billing_info(current_model,user_api_key):
    return current_model.billing_info(user_api_key)

def save_chat_history(current_model, *args):
    return current_model.save_chat_history(*args)

def save_file(filename, system, history, chatbot):
    logging.debug("会話の履歴を保存中……")
    os.makedirs(os.path.join(HISTORY_DIR), exist_ok=True)
    if filename.endswith(".json"):
        json_s = {"system": system, "history": history, "chatbot": chatbot}
        if "/" in filename or "\\" in filename:
            history_file_path = filename
        else:
            history_file_path = os.path.join(HISTORY_DIR,  filename)
        with open(history_file_path, "w") as f:
            json.dump(json_s, f)
    elif filename.endswith(".md"):
        md_s = f"system: \n- {system} \n"
        for data in history:
            md_s += f"\n{data['role']}: \n- {data['content']} \n"
        with open(os.path.join(HISTORY_DIR, filename), "w", encoding="utf8") as f:
            f.write(md_s)
    logging.debug("会話の履歴を保存しました。")
    msg = "会話の履歴を保存しました。"
    return os.path.join(HISTORY_DIR,  filename), msg

def sorted_by_pinyin(list):
    return sorted(list, key=lambda char: lazy_pinyin(char)[0][0])


def get_file_names(dir, plain=False, filetypes=[".json"]):
    #historyフォルダに保存してある会話履歴ファイル一覧の取得
    logging.debug(f"ディレクトリ{dir}，ファイルタイプは{filetypes}，純粋なテキストリスト{plain}")
    files = []
    try:
        for type in filetypes:
            files += [f for f in os.listdir(dir) if f.endswith(type)]
    except FileNotFoundError:
        files = []
    files = sorted_by_pinyin(files)
    if files == []:
        files = [""]
    logging.debug(f"files are:{files}")
    if plain:
        return files
    else:
        return gr.Dropdown.update(choices=files)

    
def get_history_names(plain=False):
    #historyフォルダに保存してある会話履歴ファイル一覧の取得
    logging.debug("履歴記録のファイル名のリストを取得します")
    return get_file_names(os.path.join(HISTORY_DIR), plain)

def load_chat_history(current_model, *args):
    return current_model.load_chat_history(*args)


def set_proxy(current_model, *args):
    current_model.set_proxy(*args)
    
    return current_model.set_proxy(*args)
    
